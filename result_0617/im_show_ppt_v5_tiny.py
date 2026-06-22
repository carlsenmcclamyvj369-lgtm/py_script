import os
import cv2
import time
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def calculate_psnr_ssim(img1, img2):
    """纯矩阵快速测算无损 PSNR & SSIM"""
    i1, i2 = img1.astype(np.float64), img2.astype(np.float64)
    mse = np.mean((i1 - i2) ** 2)
    psnr = 100.0 if mse == 0 else 20 * np.log10(255.0 / np.sqrt(mse))

    C1, C2 = (0.01 * 255) ** 2, (0.03 * 255) ** 2
    mu1, mu2 = cv2.GaussianBlur(i1, (11, 11), 1.5), cv2.GaussianBlur(i2, (11, 11), 1.5)
    sigma1_sq = cv2.GaussianBlur(i1 ** 2, (11, 11), 1.5) - mu1 ** 2
    sigma2_sq = cv2.GaussianBlur(i2 ** 2, (11, 11), 1.5) - mu2 ** 2
    sigma12 = cv2.GaussianBlur(i1 * i2, (11, 11), 1.5) - (mu1 * mu2)
    ssim_map = ((2 * mu1 * mu2 + C1) * (2 * sigma12 + C2)) / ((mu1 ** 2 + mu2 ** 2 + C1) * (sigma1_sq + sigma2_sq + C2))
    return psnr, np.mean(ssim_map)


def find_max_diff_roi(img_orig, img_denoise, roi_size=128):
    """滑动窗口自动定位灰度残差最大的 ROI"""
    h, w = img_orig.shape[:2]
    if h <= roi_size or w <= roi_size: return 0, 0, w, h
    diff_map = cv2.absdiff(cv2.cvtColor(img_orig, cv2.COLOR_BGR2GRAY), cv2.cvtColor(img_denoise, cv2.COLOR_BGR2GRAY))
    sum_map = cv2.filter2D(diff_map.astype(np.float32), -1, np.ones((roi_size, roi_size), dtype=np.float32),
                           borderType=cv2.BORDER_CONSTANT)
    _, _, _, max_loc = cv2.minMaxLoc(sum_map)
    return max(0, min(max_loc[0] - roi_size // 2, w - roi_size)), max(0, min(max_loc[1] - roi_size // 2,
                                                                             h - roi_size)), roi_size, roi_size


def format_cell(cell, text, bold=False, font_size=11, color=(255, 255, 255), bg_color=None):
    """快捷格式化 PPT 表格单元格字体与对齐方式"""
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Arial'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*bg_color)


def create_advanced_roi_dashboard(data_dir, ppt_output_path, diff_gain=10, roi_box_size=128):
    start_time = time.time()
    valid_exts = ('.png', '.jpg', '.jpeg', '.bmp', '.tiff')

    # 从 data_dir 中自动配对 mnr_input / mnr_output
    import re
    pair_pat = re.compile(r'^(.*)#(mnr_input|mnr_output)(\d+)\.bmp$')
    pairs = {}
    for fname in os.listdir(data_dir):
        if not fname.lower().endswith(valid_exts):
            continue
        m = pair_pat.match(fname)
        if not m:
            continue
        key = m.group(1) + '#' + m.group(3)  # 前缀+序号作为唯一key
        pairs.setdefault(key, {})[m.group(2)] = os.path.join(data_dir, fname)

    paired = [(k, v['mnr_input'], v['mnr_output']) for k, v in sorted(pairs.items()) if 'mnr_input' in v and 'mnr_output' in v]

    total_imgs = len(paired)
    if total_imgs == 0:
        print(f"❌ 错误: 在 {data_dir} 中未找到配对的 mnr_input/mnr_output 文件！")
        return

    temp_grid_dir = 'temp_ppt_grids'
    os.makedirs(temp_grid_dir, exist_ok=True)
    grid_images = []

    # 收集每一帧的数据，用于最后的统计图表
    metrics_summary = []

    # 1080p 标准画布参数
    canvas_w, canvas_h = 1920, 1080
    margin_top = 130
    main_col_w = (canvas_w - 70) // 2
    main_row_h = (canvas_h - margin_top - 60) // 2
    roi_col_w = (main_col_w - 20) // 2

    canvas = np.zeros((canvas_h, canvas_w, 3), dtype=np.uint8)

    print(f"\n==================== ISP IQ Board Generator ====================")
    print(f"  [+] 待处理样本总数: {total_imgs} 组")
    print(f"  [+] 输出目标报告: {ppt_output_path}")
    print(f"====================================================================")

    print("\n>>> 🚀 步骤 1: 开始计算指标并组装无损图像面板...")
    for idx, (key, in_path, out_path) in enumerate(paired):
        img_name = os.path.basename(out_path)
        img_o = cv2.imread(in_path)
        img_d = cv2.imread(out_path)
        if img_o is None or img_d is None: continue

        orig_h, orig_w = img_o.shape[:2]
        if img_o.shape != img_d.shape:
            img_d = cv2.resize(img_d, (img_o.shape[1], img_o.shape[0]), interpolation=cv2.INTER_LINEAR)

        psnr, ssim = calculate_psnr_ssim(img_o, img_d)

        # 存入列表供末尾表格使用
        metrics_summary.append({
            'id': img_name,
            'psnr': psnr,
            'ssim': ssim
        })

        print(f"  [{idx + 1} / {total_imgs}] Processing: {img_name} -> PSNR: {psnr:.2f}dB | SSIM: {ssim:.4f}")

        # ROI 寻迹与红框绘制
        rx, ry, rw, rh = find_max_diff_roi(img_o, img_d, roi_size=roi_box_size)
        roi_o = cv2.resize(img_o[ry:ry + rh, rx:rx + rw], (roi_col_w, main_row_h), interpolation=cv2.INTER_CUBIC)
        roi_d = cv2.resize(img_d[ry:ry + rh, rx:rx + rw], (roi_col_w, main_row_h), interpolation=cv2.INTER_CUBIC)

        img_o_boxed, img_d_boxed = img_o.copy(), img_d.copy()
        cv2.rectangle(img_o_boxed, (rx, ry), (rx + rw, ry + rh), (0, 0, 255), 2)
        cv2.rectangle(img_d_boxed, (rx, ry), (rx + rw, ry + rh), (0, 0, 255), 2)

        # 绝对残差计算
        img_diff = cv2.absdiff(img_o, img_d)
        img_diff = np.clip(img_diff.astype(np.float32) * diff_gain, 0, 255).astype(np.uint8)

        # 刷写画布背景
        canvas[:, :] = [0x24, 0x1E, 0x1E]

        # 顶层主客观指标文本渲染
        cv2.putText(canvas, f"| ID: {img_name}", (30, 55), cv2.FONT_HERSHEY_SIMPLEX,
                    1.25, (255, 255, 255), 3, cv2.LINE_AA)
        metrics_text = f"Metrics ->  Resolution: {orig_w}x{orig_h}  |  PSNR: {psnr:.2f} dB  |  SSIM: {ssim:.4f}  |  ROI Locate: X={rx},Y={ry}"
        cv2.putText(canvas, metrics_text, (30, 95), cv2.FONT_HERSHEY_SIMPLEX, 0.7, (0, 225, 255), 2, cv2.LINE_AA)

        # ------- 布局渲染 -------
        h_o, w_o = img_o.shape[:2]
        scale = min(main_col_w / w_o, main_row_h / h_o)
        nw, nh = int(w_o * scale), int(h_o * scale)

        left_x = 20 + (main_col_w - nw) // 2
        orig_y = margin_top + (main_row_h - nh) // 2
        deno_y = margin_top + main_row_h + 15 + (main_row_h - nh) // 2
        right_x = 40 + main_col_w + (main_col_w - nw) // 2
        diff_y = margin_top + (main_row_h - nh) // 2

        canvas[orig_y:orig_y + nh, left_x:left_x + nw] = cv2.resize(img_o_boxed, (nw, nh))
        canvas[deno_y:deno_y + nh, left_x:left_x + nw] = cv2.resize(img_d_boxed, (nw, nh))
        canvas[diff_y:diff_y + nh, right_x:right_x + nw] = cv2.resize(img_diff, (nw, nh))

        roi_y_start = margin_top + main_row_h + 15
        roi_o_x = 40 + main_col_w
        roi_d_x = roi_o_x + roi_col_w + 15
        canvas[roi_y_start:roi_y_start + main_row_h, roi_o_x:roi_o_x + roi_col_w] = roi_o
        canvas[roi_y_start:roi_y_start + main_row_h, roi_d_x:roi_d_x + roi_col_w] = roi_d

        # 局部排版文本小标签
        cv2.putText(canvas, "1. Original Full Image", (left_x, orig_y - 8), cv2.FONT_HERSHEY_SIMPLEX, 0.45,
                    (200, 200, 200), 1, cv2.LINE_AA)
        cv2.putText(canvas, "2. Denoised Full Image", (left_x, deno_y - 8), cv2.FONT_HERSHEY_SIMPLEX, 0.45,
                    (200, 200, 200), 1, cv2.LINE_AA)
        cv2.putText(canvas, f"3. Full Residual Error (Gain x{diff_gain})", (right_x, diff_y - 8),
                    cv2.FONT_HERSHEY_SIMPLEX, 0.45, (0, 165, 255), 1, cv2.LINE_AA)
        cv2.putText(canvas, "4. Original ROI (4x)", (roi_o_x, roi_y_start - 8), cv2.FONT_HERSHEY_SIMPLEX, 0.4,
                    (0, 0, 255), 1, cv2.LINE_AA)
        cv2.putText(canvas, "5. Denoised ROI (4x)", (roi_d_x, roi_y_start - 8), cv2.FONT_HERSHEY_SIMPLEX, 0.4,
                    (0, 255, 0), 1, cv2.LINE_AA)

        grid_img_path = os.path.join(temp_grid_dir, f"roi_dashboard_{idx + 1}.png")
        cv2.imwrite(grid_img_path, canvas, [int(cv2.IMWRITE_PNG_COMPRESSION), 9])
        grid_images.append(grid_img_path)

    print("\n>>> 📦 步骤 2: 正在高速向 PPT 载入对比看板页面...")
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for p_idx, img_path in enumerate(grid_images):
        print(f"  [看板载入进度: {p_idx + 1}/{len(grid_images)}]")
        slide = prs.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x24)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # 🌟 步骤 3: 统计计算 + 原生表格绘制（不含 QF）
    print("\n>>> 📊 步骤 3: 正在计算全样本平均值与中位数，自动生成末尾总结表格...")

    all_psnrs = [item['psnr'] for item in metrics_summary]
    all_ssims = [item['ssim'] for item in metrics_summary]

    mean_psnr, median_psnr = np.mean(all_psnrs), np.median(all_psnrs)
    mean_ssim, median_ssim = np.mean(all_ssims), np.median(all_ssims)

    # 设定每页表格最大行数（16行数据），支持跨页自动对齐
    rows_per_table_page = 16

    for start_i in range(0, len(metrics_summary), rows_per_table_page):
        page_items = metrics_summary[start_i: start_i + rows_per_table_page]
        is_last_table_page = (start_i + rows_per_table_page >= len(metrics_summary))

        # 1表头 + 样本数据行
        total_table_rows = 1 + len(page_items)
        if is_last_table_page:
            total_table_rows += 2  # 尾页追加 Mean 和 Median

        slide = prs.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x24)

        # 表格大标题
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Objective Metrics Benchmark Summary"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # 创建原生 3 列图元（去掉 QF 列，空间平摊更宽敞）
        left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.35 * total_table_rows)
        table_shape = slide.shapes.add_table(total_table_rows, 3, left, top, width, height)
        table = table_shape.table

        # 列宽黄金比例调节
        table.columns[0].width = Inches(6.7)  # ID
        table.columns[1].width = Inches(2.5)  # PSNR
        table.columns[2].width = Inches(2.5)  # SSIM

        # 1. 绘制深色系表头
        headers = ["Test Image ID", "PSNR (dB)", "SSIM"]
        for col_idx, h_text in enumerate(headers):
            format_cell(table.cell(0, col_idx), h_text, bold=True, font_size=12, color=(0, 225, 255),
                        bg_color=(0x2A, 0x2A, 0x32))

        # 2. 刷入当前页明细，支持斑马线错色防眼花
        for r_idx, item in enumerate(page_items):
            row_num = r_idx + 1
            bg = (0x24, 0x24, 0x2C) if row_num % 2 == 0 else (0x1E, 0x1E, 0x24)

            format_cell(table.cell(row_num, 0), item['id'], bold=False, font_size=11, color=(240, 240, 240),
                        bg_color=bg)
            format_cell(table.cell(row_num, 1), f"{item['psnr']:.3f}", bold=False, font_size=11, color=(255, 255, 255),
                        bg_color=bg)
            format_cell(table.cell(row_num, 2), f"{item['ssim']:.4f}", bold=False, font_size=11, color=(255, 255, 255),
                        bg_color=bg)

        # 3. 最后一页表格最下方，强行追加金黄色的 Mean 和绿宝石色的 Median
        if is_last_table_page:
            avg_row = total_table_rows - 2
            med_row = total_table_rows - 1

            # Average 均值行
            format_cell(table.cell(avg_row, 0), "🎯 Total Average (Mean)", bold=True, font_size=12, color=(255, 215, 0),
                        bg_color=(0x32, 0x32, 0x3A))
            format_cell(table.cell(avg_row, 1), f"{mean_psnr:.3f}", bold=True, font_size=12, color=(255, 215, 0),
                        bg_color=(0x32, 0x32, 0x3A))
            format_cell(table.cell(avg_row, 2), f"{mean_ssim:.4f}", bold=True, font_size=12, color=(255, 215, 0),
                        bg_color=(0x32, 0x32, 0x3A))

            # Median 中位数行
            format_cell(table.cell(med_row, 0), "📊 Total Median", bold=True, font_size=12, color=(0, 255, 127),
                        bg_color=(0x2A, 0x2A, 0x32))
            format_cell(table.cell(med_row, 1), f"{median_psnr:.3f}", bold=True, font_size=12, color=(0, 255, 127),
                        bg_color=(0x2A, 0x2A, 0x32))
            format_cell(table.cell(med_row, 2), f"{median_ssim:.4f}", bold=True, font_size=12, color=(0, 255, 127),
                        bg_color=(0x2A, 0x2A, 0x32))

    prs.save(ppt_output_path)

    # 清理无损大图缓存
    for img_path in grid_images:
        if os.path.exists(img_path): os.remove(img_path)
    os.rmdir(temp_grid_dir)
    print(f"\n🎉 完美收官！不含 QF 的纯净版主客观全指标 PPT 报告已生成：\n📂 {os.path.abspath(ppt_output_path)}")
    print(f"⏱️  总耗时: {time.time() - start_time:.2f} 秒")


if __name__ == '__main__':
    create_advanced_roi_dashboard(
        data_dir='out_data',
        ppt_output_path='MNR_画质对比分析报告.pptx',
        diff_gain=10,
        roi_box_size=128
    )